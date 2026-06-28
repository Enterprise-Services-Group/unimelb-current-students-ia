#!/usr/bin/env python3
"""Generate DOCX and PPTX from the fragmentation insights report."""
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN
import re

# Docx colors
NAVY_D = RGBColor(0x00, 0x0F, 0x46)
CYAN_D_D = RGBColor(0x46, 0xC8, 0xF0)
WHITE_D = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

# Pptx colors
PWHITE_D= PRGBColor(0xFF, 0xFF, 0xFF)
PNAVY_D = PRGBColor(0x00, 0x0F, 0x46)
PCYAN_D = PRGBColor(0x46, 0xC8, 0xF0)

md_path = '/Users/djmulholland/Documents/Claude/projects/unimelb-current-students-ia/report/fragmentation-insights.md'
out_dir = '/Users/djmulholland/Documents/Claude/projects/unimelb-current-students-ia/report'

# ============================================
# DOCX
# ============================================
doc = Document()

# Set default font
style = doc.styles['Normal']
font = style.font
font.name = 'Calibri'
font.size = Pt(11)
font.color.rgb = DARK_GRAY

# Parse markdown into sections
with open(md_path) as f:
    content = f.read()

# Split into sections by ## headers
sections = re.split(r'\n## ', content)

# Title
title_lines = sections[0].split('\n')
title_text = title_lines[0].replace('# ', '')
subtitle_text = title_lines[1].replace('**', '').strip() if len(title_lines) > 1 else ''
if not subtitle_text:
    sub_lines = sections[1].split('\n')[0:2] if len(sections) > 1 else ['']
    subtitle_text = sub_lines[0].replace('**', '').strip() if sub_lines else ''

title_para = doc.add_paragraph()
title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
title_run = title_para.add_run(title_text)
title_run.font.size = Pt(26)
title_run.font.color.rgb = NAVY_D
title_run.bold = True

sub_para = doc.add_paragraph()
sub_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
sub_run = sub_para.add_run(subtitle_text)
sub_run.font.size = Pt(14)
sub_run.font.color.rgb = DARK_GRAY

# Executive Summary  
es_section = next((s for s in sections if s.startswith('Executive Summary')), '')
if es_section:
    doc.add_paragraph()
    h = doc.add_heading('Executive Summary', level=2)
    for run in h.runs: run.font.color.rgb = NAVY_D
    for line in es_section.split('\n')[1:]:
        if line.startswith('**') and line.endswith('**'):
            p = doc.add_paragraph()
            r = p.add_run(line.strip('*'))
            r.bold = True
        elif line.strip():
            doc.add_paragraph(line.strip())

# Process numbered sections
for i, section in enumerate(sections[1:], 1):
    lines = section.strip().split('\n')
    header = lines[0].strip()
    
    doc.add_paragraph()
    h = doc.add_heading(header, level=2)
    for run in h.runs: run.font.color.rgb = NAVY_D
    
    for line in lines[1:]:
        line = line.strip()
        if not line: continue
        
        # Bold impact/recommendation lines
        if line.startswith('**') and ':**' in line:
            p = doc.add_paragraph()
            parts = line.split(':**', 1)
            r = p.add_run(parts[0].strip('*') + ':')
            r.bold = True
            r.font.color.rgb = NAVY_D
            if len(parts) > 1:
                p.add_run(parts[1])
        elif line.startswith('**') and line.endswith('**'):
            p = doc.add_paragraph()
            r = p.add_run(line.strip('*'))
            r.bold = True
            r.font.color.rgb = NAVY_D
        elif line.startswith('|'):
            # Table rows - add as formatted paragraph
            cells = [c.strip() for c in line.split('|')[1:-1]]
            if all(c.startswith('---') for c in cells):
                continue  # skip separator rows
            if all(c.startswith('**') for c in cells):
                continue  # skip bold header rows already captured
            p = doc.add_paragraph('  ' + '  |  '.join(cells))
            p.paragraph_format.space_before = Pt(2)
            p.paragraph_format.space_after = Pt(2)
        elif line.startswith('```'):
            continue
        else:
            doc.add_paragraph(line)

doc.save(f'{out_dir}/fragmentation-insights.docx')
print('DOCX saved')

# ============================================
# PPTX
# ============================================
prs = Presentation()
prs.slide_width = PInches(13.333)
prs.slide_height = PInches(7.5)

def add_slide(title_text, body_lines, is_title_slide=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = PRGBColor(0xFF, 0xFF, 0xFF) if not is_title_slide else PRGBColor(0x00, 0x0F, 0x46)
    
    if is_title_slide:
        # Title slide
        txBox = slide.shapes.add_textbox(PInches(1), PInches(2), PInches(11), PInches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = PPt(40)
        p.font.color.rgb = PRGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        
        if body_lines:
            txBox2 = slide.shapes.add_textbox(PInches(1), PInches(4.5), PInches(11), PInches(1))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = body_lines[0]
            p2.font.size = PPt(18)
            p2.font.color.rgb = PRGBColor(0x46, 0xC8, 0xF0)
        return
    
    # Navy bar at top
    shape = slide.shapes.add_shape(1, PInches(0), PInches(0), PInches(13.333), PInches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRGBColor(0x00, 0x0F, 0x46)
    shape.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(PInches(0.8), PInches(0.3), PInches(11.5), PInches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = PPt(28)
    p.font.color.rgb = PRGBColor(0x00, 0x0F, 0x46)
    p.font.bold = True
    
    # Body
    if body_lines:
        txBox2 = slide.shapes.add_textbox(PInches(0.8), PInches(1.4), PInches(11.5), PInches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        
        for i, line in enumerate(body_lines[:12]):  # max 12 lines per slide
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            
            line = line.strip()
            if line.startswith('**') and ':**' in line:
                parts = line.split(':**', 1)
                run = p.add_run()
                run.text = parts[0].strip('*') + ':'
                run.font.size = PPt(16)
                run.font.bold = True
                run.font.color.rgb = PRGBColor(0x00, 0x0F, 0x46)
                if len(parts) > 1:
                    run2 = p.add_run()
                    run2.text = parts[1]
                    run2.font.size = PPt(16)
            elif line.startswith('**') and line.endswith('**'):
                p.text = line.strip('*')
                p.font.size = PPt(16)
                p.font.bold = True
            elif line.startswith('|'):
                cells = [c.strip() for c in line.split('|')[1:-1]]
                if all(c.startswith('---') for c in cells):
                    p.text = ''
                    continue
                p.text = '  •  '.join(cells[:4])
                p.font.size = PPt(14)
            elif line.startswith('```'):
                continue
            else:
                p.text = line
                p.font.size = PPt(16)
            
            p.space_after = PPt(6)

# Title slide
add_slide("Fragmentation, Duplication\n& the Student Experience", 
          ["Insights from the University of Melbourne Web Estate Audit  •  June 2026"], 
          is_title_slide=True)

# Build slides from sections
current_section = ''
current_lines = []

def flush_section():
    global current_section, current_lines
    if current_section and current_lines:
        # Split long sections across multiple slides
        if len(current_lines) > 10:
            add_slide(current_section, current_lines[:10])
            remaining = current_lines[10:]
            while remaining:
                add_slide(f"{current_section} (cont.)", remaining[:10])
                remaining = remaining[10:]
        else:
            add_slide(current_section, current_lines)
    current_section = ''
    current_lines = []

# Process content
for i, section in enumerate(sections[1:]):
    lines = section.strip().split('\n')
    header = lines[0].strip()
    
    # Flush previous
    flush_section()
    
    current_section = header
    for line in lines[1:]:
        line = line.strip()
        if line and not line.startswith('#'):
            current_lines.append(line)

flush_section()

# Recommendations slide
rec_lines = [
    "Immediate (this semester):",
    "1. Standardise URL convention — every faculty at /students",
    "2. Add Prospective → Current → Alumni lifecycle navigation",
    "3. Build graduation pathway page on students.unimelb",
    "",
    "Short-term (next 6 months):",
    "4. Consolidate 70 link-farm pages into central hub",
    "5. Wire scholarships.unimelb as canonical source",
    "6. Move Moving Guide, GDP, Glossary from study to students",
    "7. Consolidate MDHS/FBE school-level CS sections",
    "8. Create cross-faculty topic aggregation pages",
    "",
    "Medium-term (12 months):",
    "9. University-wide minimum content standards",
    "10. Build prospective→enrolled recognition pathway",
    "11. Single signposting framework for all high-stakes questions",
]
add_slide("12 Recommendations", rec_lines)

# Data slide
data_lines = [
    "36,500+ pages across 19 University websites — full HTML crawl, June 2026",
    "",
    "65% of faculty content is genuinely unique (not link-farm)",
    "22 topics fragmented across 9–11 faculty sites each",
    "5 different URL conventions for 'current student information'",
    "Zero links from current-student estate to alumni site",
    "3-way scholarship split across dedicated domain + faculty + study sites",
    "Average faculty site: 2,100 pages (range: 83–6,079)",
    "",
    "All findings evidence-grounded and traceable to specific pages",
]
add_slide("About the Data", data_lines)

prs.save(f'{out_dir}/fragmentation-insights.pptx')
print('PPTX saved')
print('Done — MD, DOCX, PPTX complete')
