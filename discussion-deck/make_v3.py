#!/usr/bin/env python3
"""Build v3 from the hand-edited v2 .pptx (NOT from build-v2.js).

The v2 .pptx has manual visual edits that are not in build-v2.js, so v3 is
created by COPYING v2 and ADDING slides via python-pptx — existing slides are
left byte-for-byte untouched. Adds three slides reflecting the full-faculty +
cross-corpus lifecycle analysis:
  A (main flow, after 'Three phases'): four websites, two broken seams
  B (main flow):                       inside the faculty estate / research paradox
  C (appendix):                        cross-corpus evidence table
Then renumbers the static footer page numbers (text-only, no layout change).

Re-runnable: always starts from a fresh copy of v2. Run from repo root or here:
  python3 discussion-deck/make_v3.py
"""
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).resolve().parent
SRC = HERE / "Current-Students-EndToEnd-Discussion-v2.pptx"
DST = HERE / "Current-Students-EndToEnd-Discussion-v3.pptx"

# ---- deck palette + fonts (matched to build-v2.js / the existing deck) ----
def col(h): return RGBColor.from_string(h)
NAVY, NAVY2, CYAN, CYANDK = "000F46", "000B34", "46C8F0", "2F95B7"
INK, LIGHT, TINT, MUTED = "1C1C28", "F4F4F4", "EAF6FB", "5B5B6B"
WHITE, ICE, DIM, LINELT, FAINT = "FFFFFF", "C9CDE6", "AAB0D4", "DDDDE7", "8F96C2"
HEAD, BODY, MONO = "Fraunces", "Source Sans 3", "Source Code Pro"

shutil.copy(SRC, DST)
prs = Presentation(str(DST))


# ---------------------------------------------------------------- helpers ---
def add_slide():
    s = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in list(s.placeholders):           # strip inherited layout placeholders
        ph._element.getparent().remove(ph._element)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = col(WHITE)
    return s


def box(s, x, y, w, h, fill=None, line=None, lw=0.5):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = col(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = col(line); sh.line.width = Pt(lw)
    if sh.has_text_frame:
        sh.text_frame.word_wrap = True
    return sh


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, sp=None):
    """runs: list of dicts or list-of-lists-of-dicts (paragraphs). dict keys:
       t,font,size,color,bold,italic,spacing(charSpacing pt*100)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if runs and isinstance(runs[0], list) else [runs]
    for pi, prun in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        if sp: p.line_spacing = sp
        for r in prun:
            run = p.add_run(); run.text = r["t"]
            f = run.font
            f.name = r.get("font", BODY); f.size = Pt(r.get("size", 12))
            f.bold = r.get("bold", False); f.italic = r.get("italic", False)
            f.color.rgb = col(r.get("color", INK))
    return tb


def head(s, title, page):
    box(s, 0.6, 0.62, 0.28, 0.28, fill=CYAN)
    txt(s, 1.0, 0.5, 11.7, 0.55, [{"t": title, "font": HEAD, "size": 29, "color": NAVY, "bold": True}],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 0.6, 7.04, 10.5, 0.3,
        [{"t": "University of Melbourne — web estate alignment discussion · for discussion",
          "size": 9, "color": MUTED}])
    txt(s, 12.3, 7.04, 0.45, 0.3, [{"t": str(page), "size": 9, "color": MUTED}], align=PP_ALIGN.RIGHT)


def set_notes(s, t):
    tf = s.notes_slide.notes_text_frame
    if tf is not None:
        tf.text = t


def discuss(s, y, prompt):
    box(s, 0.62, y, 11.9, 0.68, fill=TINT)
    box(s, 0.62, y, 0.1, 0.68, fill=CYAN)
    txt(s, 0.95, y, 11.3, 0.68,
        [{"t": "To discuss:  ", "color": NAVY, "bold": True, "size": 13},
         {"t": prompt, "color": NAVY2, "size": 13}],
        anchor=MSO_ANCHOR.MIDDLE)


# ============================================================= SLIDE A =====
# four websites, two broken seams
A = add_slide(); head(A, "The lifecycle: four websites, two broken seams", 0)
txt(A, 0.62, 1.28, 11.9, 0.8,
    [{"t": "Beyond the three phases: the journey actually crosses four separate websites — each with its own "
           "architecture and vocabulary. Two of the three handoffs have no web-level bridge at all.",
      "size": 15, "color": INK}], sp=1.12)

# 4 domain boxes (deck phase colours: tint=before, navy2/navy=during, light=after)
bx = [(0.62, "study.unimelb", "prospective · 300pp", TINT, CYANDK, NAVY),
      (3.82, "9 faculty sites", "current · ~1,000pp", NAVY2, None, WHITE),
      (7.02, "students.unimelb", "current · central", NAVY, None, WHITE),
      (10.22, "alumni site", "alumni · 116pp", LIGHT, LINELT, NAVY)]
BW, BY, BH = 2.5, 2.5, 1.05
for x, t1, t2, fill, ln, tc in bx:
    box(A, x, BY, BW, BH, fill=fill, line=ln, lw=1.25)
    txt(A, x, BY + 0.22, BW, 0.4, [{"t": t1, "font": BODY, "size": 13, "color": tc, "bold": True}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(A, x, BY + 0.62, BW, 0.3, [{"t": t2, "font": MONO, "size": 9, "color": tc}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# seam glyphs in the gaps (mid box height)
gy = BY + BH / 2 - 0.18
for gx, glyph, gc in [(3.12, "✗", "A32D2D"), (6.32, "↔", CYANDK), (9.52, "✗", "A32D2D")]:
    txt(A, gx, gy, 0.7, 0.36, [{"t": glyph, "font": BODY, "size": 18, "color": gc, "bold": True}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(A, 0.62, BY + BH + 0.08, 11.9, 0.3,
    [{"t": "✗  broken handoff — no shared vocabulary, no web-level bridge        ↔  linked (same stage)",
      "font": MONO, "size": 9.5, "italic": True, "color": MUTED}], align=PP_ALIGN.CENTER)

# 3 stat chips
chips = [("RESEARCH PARADOX", "MDHS: 63 research pages, 1 HDR-support page"),
         ("CAREERS ACROSS LIFE", "Authored in ~10 places across 4 domains, 0 cross-links"),
         ("SMALLEST TENANT", "Current students = 4% of the faculty web estate")]
CY, CW, CH = 4.15, 3.9, 1.05
for i, (t1, t2) in enumerate(chips):
    x = 0.62 + i * (CW + 0.2)
    box(A, x, CY, CW, CH, fill=TINT)
    box(A, x, CY, 0.08, CH, fill=CYAN)
    txt(A, x + 0.22, CY + 0.16, CW - 0.34, 0.3,
        [{"t": t1, "font": MONO, "size": 10, "color": CYANDK, "bold": True, "spacing": 1}])
    txt(A, x + 0.22, CY + 0.5, CW - 0.34, 0.5, [{"t": t2, "size": 11.5, "color": NAVY2}], sp=1.08)

discuss(A, 5.55, "Which handoff matters most to fix first — apply→enrol, or graduate→alumni? "
                 "And where would a shared vocabulary help users cross between sites?")
set_notes(A,
    "Deepens the 'Three phases' slide (~4 min). The journey crosses FOUR separate websites, not one. Two of the "
    "three handoffs have no web link at all: apply→enrol (study.unimelb to the current estate, no shared "
    "vocabulary) and graduate→alumni (0 links, confirmed in the crawl). The three chips are the headline "
    "cross-corpus findings. Source: analysis/lifecycle-cross-corpus.md.")

# ============================================================= SLIDE B =====
# inside the faculty estate / research paradox
B = add_slide(); head(B, "Inside the faculty estate — the 4% tenant", 0)
txt(B, 0.62, 1.28, 6.4, 4.6,
    [[{"t": "A crawl of the full faculty domains (2,363 pages across 9 sites) shows current students are the "
            "smallest named tenant:", "size": 15, "color": INK}],
     [{"t": "", "size": 6}],
     [{"t": "about 19%  ·  news 11%  ·  research 9%  ·  prospective study 8%  ·  current students 4%",
       "size": 13, "color": NAVY2, "bold": True}],
     [{"t": "", "size": 6}],
     [{"t": "The faculty sites are built for reputation and recruitment — a current student is a minority "
            "audience on the very domains that host their faculty's content.", "size": 15, "color": INK}],
     [{"t": "", "size": 6}],
     [{"t": "And research reputation is decoupled from research-student support: the faculties with the largest "
            "research estates carry the least HDR-candidature support  →", "size": 15, "color": INK}]], sp=1.16)

# research paradox box-grid (right)
GX, GW = 7.35, 5.37
rows = [("Faculty", "Research", "HDR support", True),
        ("MDHS", "63", "1", False), ("Science", "60", "6", False),
        ("Education", "44", "0", False), ("ABP / MSD", "46", "29", False),
        ("Law", "19", "79 *", False)]
ry, rh = 1.55, 0.62
for i, (a, b, c, hdr) in enumerate(rows):
    y = ry + i * rh
    box(B, GX, y, GW, rh, fill=(NAVY if hdr else WHITE), line=LINELT, lw=1)
    fc = WHITE if hdr else NAVY2
    fb = True if hdr else False
    txt(B, GX + 0.15, y, 2.6, rh, [{"t": a, "size": 11.5, "color": fc, "bold": (fb or not hdr)}],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(B, GX + 2.7, y, 1.3, rh, [{"t": b, "font": MONO, "size": 12, "color": fc, "bold": fb}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(B, GX + 4.0, y, 1.25, rh, [{"t": c, "font": MONO, "size": 12, "color": fc, "bold": fb}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(B, GX, ry + len(rows) * rh + 0.06, GW, 0.5,
    [{"t": "* Law's 79 is mostly researcher-profile pages (61 bios); only ~18 are genuine candidature support.",
      "font": BODY, "size": 10, "italic": True, "color": MUTED}], sp=1.05)

discuss(B, 5.95, "Does the faculty web presence serve current students, or are they an afterthought on a "
                 "reputation-and-recruitment site? Who owns the research-student experience?")
set_notes(B,
    "The faculty-estate context (~4 min). Two points: (1) current students are only ~4-7% of each faculty's "
    "public site — the rest is reputation/recruitment. (2) The research paradox: the most research-intensive "
    "faculties (MDHS, Science, Education) provide the least HDR-student support; Law looks supported only "
    "because 61 of its 79 pages are researcher bios. Source: analysis/full-faculty-context.md + "
    "lifecycle-cross-corpus.md.")

# ============================================================= SLIDE C =====
# appendix: cross-corpus evidence table
Cs = add_slide(); head(Cs, "Cross-corpus evidence — services across the lifecycle", 0)
txt(Cs, 0.62, 1.28, 11.9, 0.6,
    [{"t": "Pages per topic at each lifecycle stage (June 2026 crawl). Cross-cutting services are authored at "
           "every stage; placements, course planning and candidature are genuinely stage-bound.",
      "size": 13.5, "color": INK}], sp=1.12)

data = [("Topic / service", "Prospective", "Current", "Alumni", "Pattern"),
        ("Careers & employability", "38", "124", "38", "all three stages"),
        ("Scholarships", "75", "116", "1", "prospective + current"),
        ("Student life", "156", "108", "21", "all three stages"),
        ("International", "121", "100", "20", "all three stages"),
        ("Wellbeing & health", "17", "46", "17", "all three stages"),
        ("Research candidature", "1", "192", "1", "current only"),
        ("Placements & WIL", "10", "230", "0", "current only (faculty-owned)")]
tx, ty, tw = 0.62, 2.15, 12.1
colw = [3.5, 1.9, 1.7, 1.5, 3.5]
rh = 0.5
gf = Cs.shapes.add_table(len(data), 5, Inches(tx), Inches(ty), Inches(tw), Inches(rh * len(data)))
tbl = gf.table
tbl.first_row = False; tbl.horz_banding = False
for ci, w in enumerate(colw):
    tbl.columns[ci].width = Inches(w)
for ri, row in enumerate(data):
    tbl.rows[ri].height = Inches(rh)
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = col(NAVY if ri == 0 else WHITE)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci in (0, 4) else PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.name = (BODY if ci == 0 or ci == 4 else MONO)
        r.font.size = Pt(11)
        r.font.bold = (ri == 0) or (ci == 0)
        r.font.color.rgb = col(WHITE if ri == 0 else (NAVY2 if ci == 0 else INK))
txt(Cs, 0.62, ty + rh * len(data) + 0.12, 11.9, 0.3,
    [{"t": "Source: analysis/lifecycle-cross-corpus.md · data/clean/cross_corpus_lifecycle.py "
           "(reproducible). Careers is additionally authored at ~10 distinct locations across 4 domains.",
      "font": MONO, "size": 9, "italic": True, "color": MUTED}])
set_notes(Cs,
    "Reference only. Topic pages at each lifecycle stage. Cross-cutting services (careers, student life, "
    "international, wellbeing) recur at all three stages — each authored in three voices for the same person at "
    "different ages. Stage-bound content (placements, course planning, candidature) is what stays with "
    "faculties. Full method: data/clean/cross_corpus_lifecycle.py.")

# ---- reorder: A,B into main flow after 'Three phases' (orig index 4) -------
# original slides 0..12; appended A=13, B=14, C=15. Target: A->5, B->6, C stays last.
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
A_el, B_el = ids[13], ids[14]            # C_el = ids[15] stays at end
sldIdLst.remove(A_el); sldIdLst.insert(5, A_el)
sldIdLst.remove(B_el); sldIdLst.insert(6, B_el)

# ---- renumber static footer page numbers (text-only) -----------------------
for n, s in enumerate(prs.slides, start=1):
    for sh in s.shapes:
        if not sh.has_text_frame or sh.left is None:
            continue
        # footer page-number box: far right, bottom
        if sh.left > Inches(12.0) and sh.top > Inches(6.9):
            run = sh.text_frame.paragraphs[0].runs
            if run and run[0].text.strip().isdigit():
                run[0].text = str(n)
    # title slide (n==1) has no footer number; skip silently

prs.save(str(DST))
print(f"wrote {DST.name} · {len(prs.slides._sldIdLst)} slides")
